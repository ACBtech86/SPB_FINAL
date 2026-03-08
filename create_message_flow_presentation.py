#!/usr/bin/env python3
"""
Create animated PowerPoint presentation showing SPB message flow
through IBM MQ, PostgreSQL, BACEN Auto Responder, and SPBSite
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# Define colors
COLOR_TITLE = RGBColor(0, 51, 102)  # Dark blue
COLOR_HEADER = RGBColor(68, 114, 196)  # Medium blue
COLOR_MQ = RGBColor(255, 192, 0)  # Yellow
COLOR_DB = RGBColor(112, 173, 71)  # Green
COLOR_BACEN = RGBColor(237, 125, 49)  # Orange
COLOR_SPBSITE = RGBColor(91, 155, 213)  # Light blue
COLOR_TEST = RGBColor(192, 0, 0)  # Red
COLOR_ARROW = RGBColor(0, 112, 192)  # Blue

def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue background

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "SPB Message Flow Integration"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_text = subtitle_frame.add_paragraph()
    subtitle_text.text = "Complete Integration Test with Automated BACEN Response"
    subtitle_text.font.size = Pt(32)
    subtitle_text.font.color.rgb = COLOR_HEADER
    subtitle_text.alignment = PP_ALIGN.CENTER

    # Components
    components_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1.5))
    components_frame = components_box.text_frame
    components_text = components_frame.add_paragraph()
    components_text.text = "IBM MQ • PostgreSQL • BACEN Auto Responder • SPBSite"
    components_text.font.size = Pt(24)
    components_text.font.color.rgb = RGBColor(100, 100, 100)
    components_text.alignment = PP_ALIGN.CENTER

def add_architecture_slide(prs):
    """Slide 2: Architecture overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Complete Architecture Overview"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE

    # Test Component (Top)
    test_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(1), Inches(4.5), Inches(0.8)
    )
    test_box.fill.solid()
    test_box.fill.fore_color.rgb = COLOR_TEST
    test_box.line.color.rgb = RGBColor(0, 0, 0)
    test_frame = test_box.text_frame
    test_frame.text = "Message Flow Test\n(test_message_flow.py)"
    test_frame.paragraphs[0].font.size = Pt(18)
    test_frame.paragraphs[0].font.bold = True
    test_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    test_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    test_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # SPBSite
    spbsite_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(2.5), Inches(3), Inches(1.5)
    )
    spbsite_box.fill.solid()
    spbsite_box.fill.fore_color.rgb = COLOR_SPBSITE
    spbsite_box.line.color.rgb = RGBColor(0, 0, 0)
    spbsite_frame = spbsite_box.text_frame
    spbsite_frame.text = "SPBSite Web Interface\nPort: 8000\n\n• Monitoring Pages\n• Message Viewer\n• API Docs"
    spbsite_frame.paragraphs[0].font.size = Pt(16)
    spbsite_frame.paragraphs[0].font.bold = True
    spbsite_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # BACEN Auto Responder
    bacen_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(2.5), Inches(3), Inches(1.5)
    )
    bacen_box.fill.solid()
    bacen_box.fill.fore_color.rgb = COLOR_BACEN
    bacen_box.line.color.rgb = RGBColor(0, 0, 0)
    bacen_frame = bacen_box.text_frame
    bacen_frame.text = "BACEN Auto Responder\nPoll: 0.5s\n\n• Auto Detect\n• Generate Response\n• UTF-16BE Encoding"
    bacen_frame.paragraphs[0].font.size = Pt(16)
    bacen_frame.paragraphs[0].font.bold = True
    bacen_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # IBM MQ
    mq_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(2.5), Inches(3), Inches(1.5)
    )
    mq_box.fill.solid()
    mq_box.fill.fore_color.rgb = COLOR_MQ
    mq_box.line.color.rgb = RGBColor(0, 0, 0)
    mq_frame = mq_box.text_frame
    mq_frame.text = "IBM MQ 9.3.0.0\nQM.36266751.01\n\n• IF Staging Queues\n• BACEN Inbound\n• 14 Queues Total"
    mq_frame.paragraphs[0].font.size = Pt(16)
    mq_frame.paragraphs[0].font.bold = True

    # PostgreSQL
    db_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(5), Inches(5.5), Inches(1.5)
    )
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = COLOR_DB
    db_box.line.color.rgb = RGBColor(0, 0, 0)
    db_frame = db_box.text_frame
    db_frame.text = "PostgreSQL Database (bcspbstr)\n\n• SPB_BACEN_TO_LOCAL • SPB_LOCAL_TO_BACEN • SPB_LOG_BACEN\n• 979 Message Types • 581 Field Types • 32,955 Fields"
    db_frame.paragraphs[0].font.size = Pt(14)
    db_frame.paragraphs[0].font.bold = True
    db_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Add arrows
    # Test to SPBSite
    slide.shapes.add_connector(1, Inches(5), Inches(1.8), Inches(2), Inches(2.5))
    # Test to BACEN
    slide.shapes.add_connector(1, Inches(6.75), Inches(1.8), Inches(5.5), Inches(2.5))
    # Test to MQ
    slide.shapes.add_connector(1, Inches(8.5), Inches(1.8), Inches(9), Inches(2.5))
    # BACEN to MQ
    slide.shapes.add_connector(1, Inches(7), Inches(3.25), Inches(7.5), Inches(3.25))
    # All to DB
    slide.shapes.add_connector(1, Inches(2), Inches(4), Inches(6.75), Inches(5))
    slide.shapes.add_connector(1, Inches(5.5), Inches(4), Inches(6.75), Inches(5))
    slide.shapes.add_connector(1, Inches(9), Inches(4), Inches(6.75), Inches(5))

def add_flow_step_1(prs):
    """Slide 3: Step 1 - Send message to IF staging queue"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Step 1: Send Message to IF Staging Queue"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE

    # Test sends message
    test_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5), Inches(3), Inches(1)
    )
    test_box.fill.solid()
    test_box.fill.fore_color.rgb = COLOR_TEST
    test_box.line.color.rgb = RGBColor(0, 0, 0)
    test_frame = test_box.text_frame
    test_frame.text = "Test Script\nSends GEN0001"
    test_frame.paragraphs[0].font.size = Pt(18)
    test_frame.paragraphs[0].font.bold = True
    test_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    test_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    test_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow
    arrow = slide.shapes.add_connector(1, Inches(4), Inches(2), Inches(6.5), Inches(2))
    arrow.line.color.rgb = COLOR_ARROW
    arrow.line.width = Pt(3)

    # IF Staging Queue
    queue_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(1.5), Inches(5), Inches(1)
    )
    queue_box.fill.solid()
    queue_box.fill.fore_color.rgb = COLOR_MQ
    queue_box.line.color.rgb = RGBColor(0, 0, 0)
    queue_frame = queue_box.text_frame
    queue_frame.text = "IF Staging Queue\nQL.36266751.01.ENTRADA.IF"
    queue_frame.paragraphs[0].font.size = Pt(18)
    queue_frame.paragraphs[0].font.bold = True
    queue_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    queue_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Message details
    details_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(10.5), Inches(3))
    details_frame = details_box.text_frame
    details_text = """Message Details:
• Message Type: GEN0001 (Echo Request)
• Sender: BACEN (00038166)
• Receiver: FINVEST (36266751)
• Encoding: UTF-16BE
• Queue: IF Staging (Finvest Outbound)
• Status: ✅ Message queued successfully"""
    details_frame.text = details_text
    details_frame.paragraphs[0].font.size = Pt(20)
    details_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

def add_flow_step_2(prs):
    """Slide 4: Step 2 - BACEN Auto Responder processes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Step 2: BACEN Auto Responder Processes Message"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE

    # IF Staging Queue
    queue_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5), Inches(3.5), Inches(0.8)
    )
    queue_box.fill.solid()
    queue_box.fill.fore_color.rgb = COLOR_MQ
    queue_box.line.color.rgb = RGBColor(0, 0, 0)
    queue_frame = queue_box.text_frame
    queue_frame.text = "IF Staging Queue"
    queue_frame.paragraphs[0].font.size = Pt(16)
    queue_frame.paragraphs[0].font.bold = True
    queue_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    queue_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow to BACEN
    arrow1 = slide.shapes.add_connector(1, Inches(4.5), Inches(1.9), Inches(5.5), Inches(1.9))
    arrow1.line.color.rgb = COLOR_ARROW
    arrow1.line.width = Pt(3)

    # BACEN Auto Responder
    bacen_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(1.2), Inches(4), Inches(1.5)
    )
    bacen_box.fill.solid()
    bacen_box.fill.fore_color.rgb = COLOR_BACEN
    bacen_box.line.color.rgb = RGBColor(0, 0, 0)
    bacen_frame = bacen_box.text_frame
    bacen_frame.text = "BACEN Auto Responder\n\n1. Detect message\n2. Decode UTF-16BE\n3. Parse XML\n4. Generate GEN0002"
    bacen_frame.paragraphs[0].font.size = Pt(14)
    bacen_frame.paragraphs[0].font.bold = True
    bacen_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Arrow to BACEN Inbound
    arrow2 = slide.shapes.add_connector(1, Inches(9.5), Inches(1.9), Inches(10.5), Inches(1.9))
    arrow2.line.color.rgb = COLOR_ARROW
    arrow2.line.width = Pt(3)

    # BACEN Inbound Queue
    inbound_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(10.5), Inches(1.5), Inches(2), Inches(0.8)
    )
    inbound_box.fill.solid()
    inbound_box.fill.fore_color.rgb = COLOR_MQ
    inbound_box.line.color.rgb = RGBColor(0, 0, 0)
    inbound_frame = inbound_box.text_frame
    inbound_frame.text = "BACEN Inbound"
    inbound_frame.paragraphs[0].font.size = Pt(14)
    inbound_frame.paragraphs[0].font.bold = True
    inbound_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    inbound_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Processing details
    details_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.5), Inches(3))
    details_frame = details_box.text_frame
    details_text = """Processing Details:
• Poll Interval: 0.5 seconds
• Detected: GEN0001 request in IF staging queue
• Decoded: UTF-16BE with byte-swap → UTF-8
• Parsed: Extracted NUOp, CodMsg, Emissor, Destinatário
• Generated: GEN0002 response with success status (00)
• Encoded: UTF-8 → UTF-16BE with byte-swap
• Sent: Response to QL.REQ.00038166.36266751.01
• Correlation ID: Matched to original message
• Time: < 2 seconds (automated)"""
    details_frame.text = details_text
    details_frame.paragraphs[0].font.size = Pt(16)
    details_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

def add_flow_step_3(prs):
    """Slide 5: Step 3 - Database storage"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Step 3: Database Storage & SPBSite Display"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE

    # BACEN Inbound Queue
    queue_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5), Inches(3), Inches(0.8)
    )
    queue_box.fill.solid()
    queue_box.fill.fore_color.rgb = COLOR_MQ
    queue_box.line.color.rgb = RGBColor(0, 0, 0)
    queue_frame = queue_box.text_frame
    queue_frame.text = "BACEN Inbound Queue"
    queue_frame.paragraphs[0].font.size = Pt(14)
    queue_frame.paragraphs[0].font.bold = True
    queue_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    queue_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow to Test
    arrow1 = slide.shapes.add_connector(1, Inches(4), Inches(1.9), Inches(5), Inches(1.9))
    arrow1.line.color.rgb = COLOR_ARROW
    arrow1.line.width = Pt(3)

    # Test retrieves
    test_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5), Inches(1.5), Inches(2.5), Inches(0.8)
    )
    test_box.fill.solid()
    test_box.fill.fore_color.rgb = COLOR_TEST
    test_box.line.color.rgb = RGBColor(0, 0, 0)
    test_frame = test_box.text_frame
    test_frame.text = "Test Retrieves"
    test_frame.paragraphs[0].font.size = Pt(14)
    test_frame.paragraphs[0].font.bold = True
    test_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    test_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    test_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow to DB
    arrow2 = slide.shapes.add_connector(1, Inches(7.5), Inches(1.9), Inches(8.5), Inches(1.9))
    arrow2.line.color.rgb = COLOR_ARROW
    arrow2.line.width = Pt(3)

    # PostgreSQL
    db_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(8.5), Inches(1.2), Inches(4), Inches(1.5)
    )
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = COLOR_DB
    db_box.line.color.rgb = RGBColor(0, 0, 0)
    db_frame = db_box.text_frame
    db_frame.text = "PostgreSQL (bcspbstr)\n\nSPB_LOG_BACEN\nSPB_BACEN_TO_LOCAL"
    db_frame.paragraphs[0].font.size = Pt(16)
    db_frame.paragraphs[0].font.bold = True
    db_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    db_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # SPBSite
    spbsite_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(3.5), Inches(7.5), Inches(1.2)
    )
    spbsite_box.fill.solid()
    spbsite_box.fill.fore_color.rgb = COLOR_SPBSITE
    spbsite_box.line.color.rgb = RGBColor(0, 0, 0)
    spbsite_frame = spbsite_box.text_frame
    spbsite_frame.text = "SPBSite Web Interface\nhttp://localhost:8000/monitoring/messages/inbound/bacen"
    spbsite_frame.paragraphs[0].font.size = Pt(18)
    spbsite_frame.paragraphs[0].font.bold = True
    spbsite_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    spbsite_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    spbsite_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow DB to SPBSite
    arrow3 = slide.shapes.add_connector(1, Inches(10.5), Inches(2.7), Inches(6.75), Inches(3.5))
    arrow3.line.color.rgb = COLOR_ARROW
    arrow3.line.width = Pt(3)

    # Details
    details_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.5), Inches(1.8))
    details_frame = details_box.text_frame
    details_text = """Storage & Display:
• SPB_LOG_BACEN: Transaction log with complete message XML
• SPB_BACEN_TO_LOCAL: Application table (Status: P, Flag: N)
• SPBSite: Real-time web monitoring - view messages through browser
• Access: http://localhost:8000 for complete message details"""
    details_frame.text = details_text
    details_frame.paragraphs[0].font.size = Pt(16)
    details_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

def add_results_slide(prs):
    """Slide 6: Test results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Integration Test Results - All Systems Passing ✅"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE

    # Test results box
    results_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(4.5))
    results_frame = results_box.text_frame
    results_text = """Test Summary (7/7 Passing):

✅ Test 1: Send Message to IF Staging Queue - PASSED
✅ Test 2: Retrieve BACEN Auto-Response - PASSED
✅ Test 3: Log BACEN Response to Database - PASSED
✅ Test 4: Store BACEN Response in App Table - PASSED
✅ Test 5: Verify BACEN Auto Responder - PASSED
✅ Test 6: Verify SPBSite Integration - PASSED
✅ Test 7: Verify Complete Flow - PASSED

Statistics:
• Messages Processed: 1
• Responses Sent: 1
• Errors: 0
• Duration: ~10 seconds
• Automation Level: 100% (no manual intervention)"""
    results_frame.text = results_text
    results_frame.paragraphs[0].font.size = Pt(18)
    results_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)

    # Success banner
    banner_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3), Inches(6.2), Inches(7.5), Inches(0.8)
    )
    banner_box.fill.solid()
    banner_box.fill.fore_color.rgb = RGBColor(0, 176, 80)
    banner_box.line.color.rgb = RGBColor(0, 0, 0)
    banner_frame = banner_box.text_frame
    banner_frame.text = "🎉 All Tests Passed! Complete Integration Working! 🎉"
    banner_frame.paragraphs[0].font.size = Pt(24)
    banner_frame.paragraphs[0].font.bold = True
    banner_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    banner_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    banner_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_components_slide(prs):
    """Slide 7: Component details"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Integrated Components Details"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_TITLE

    # Create 4 boxes for components
    components = [
        {
            'title': 'IBM MQ 9.3.0.0',
            'details': '• Queue Manager: QM.36266751.01\n• Channel: FINVEST.SVRCONN\n• Port: 1414\n• Queues: 14 (IF + BACEN)\n• Auto-start: systemd',
            'color': COLOR_MQ,
            'x': 0.5,
            'y': 1.5
        },
        {
            'title': 'PostgreSQL 16.13',
            'details': '• Database: bcspbstr\n• Tables: 7\n• Messages: 979 types\n• Fields: 32,955 definitions\n• Drivers: psycopg2, asyncpg',
            'color': COLOR_DB,
            'x': 6.75,
            'y': 1.5
        },
        {
            'title': 'BACEN Auto Responder',
            'details': '• Auto message processing\n• UTF-16BE encoding\n• Poll interval: 0.5s\n• Response generation\n• Statistics tracking',
            'color': COLOR_BACEN,
            'x': 0.5,
            'y': 4.2
        },
        {
            'title': 'SPBSite Web Interface',
            'details': '• FastAPI + Uvicorn\n• Port: 8000\n• Monitoring pages\n• Message viewer\n• API documentation',
            'color': COLOR_SPBSITE,
            'x': 6.75,
            'y': 4.2
        }
    ]

    for comp in components:
        comp_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(comp['x']), Inches(comp['y']), Inches(6), Inches(2.2)
        )
        comp_box.fill.solid()
        comp_box.fill.fore_color.rgb = comp['color']
        comp_box.line.color.rgb = RGBColor(0, 0, 0)
        comp_frame = comp_box.text_frame
        comp_frame.text = f"{comp['title']}\n\n{comp['details']}"
        comp_frame.paragraphs[0].font.size = Pt(18)
        comp_frame.paragraphs[0].font.bold = True
        if comp['color'] == COLOR_MQ:
            comp_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        else:
            comp_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Create all slides
print("Creating PowerPoint presentation...")
add_title_slide(prs)
add_architecture_slide(prs)
add_flow_step_1(prs)
add_flow_step_2(prs)
add_flow_step_3(prs)
add_results_slide(prs)
add_components_slide(prs)

# Save presentation
output_file = '/home/ubuntu/SPB_FINAL/SPB_Message_Flow_Integration.pptx'
prs.save(output_file)
print(f"✅ PowerPoint created: {output_file}")
print(f"   Slides: 7")
print(f"   Topics: Architecture, Message Flow, Components, Test Results")
